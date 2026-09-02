#!/usr/bin/env python3
"""Генерация бизнес-плана скалодрома в формате PPTX по образцу."""

from __future__ import annotations

import math
import os
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "presentation-assets"
OUTPUT = ROOT / "Бизнес-план_скалодрома.pptx"
TEMPLATE_WAVE = (
    Path.home()
    / ".cursor/projects/Users-maksimizrailev-Documents-hexlet-studying-ski-guide-Italy"
    / "assets/image-6c8ebd01-2bc4-434a-ab85-eb57bd1a4b51.png"
)

# Реальное объявление Avito (проверено 26.05.2026)
AVITO_URL = (
    "https://www.avito.ru/sankt-peterburg/kommercheskaya_nedvizhimost/"
    "svobodnogo_naznacheniya_proizvodstvo_850_m_7985954480"
)
ADDRESS = "г. Санкт-Петербург, Кондратьевский пр-т, 2П"
AREA_M2 = 850
RENT_MONTH = 765_000
RENT_M2 = 900
CEILING_M = 5.6
DEPOSIT = 1_530_000

ORANGE = RGBColor(0xF5, 0x7C, 0x00)
ORANGE_DARK = RGBColor(0xC4, 0x3A, 0x00)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x2E, 0x7D, 0x32)

STARTUP_PREMISES = 5_160_000
STARTUP_INVENTORY = 12_570_000
STARTUP_TOTAL = STARTUP_PREMISES + STARTUP_INVENTORY
MONTHLY_TOTAL = 1_596_000
MONTHLY_INCOME = 1_562_000


def _font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/System/Library/Fonts/Helvetica.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for path in candidates:
        if os.path.exists(path):
            try:
                return ImageFont.truetype(path, size)
            except OSError:
                continue
    return ImageFont.load_default()


def _save(img: Image.Image, name: str) -> Path:
    ASSETS.mkdir(parents=True, exist_ok=True)
    path = ASSETS / name
    img.save(path, quality=92)
    return path


def make_climbing_photo(seed: int, title: str) -> Path:
    w, h = 1200, 800
    img = Image.new("RGB", (w, h), (28, 32, 42))
    draw = ImageDraw.Draw(img)

    # пол / мат
    draw.rectangle([0, int(h * 0.72), w, h], fill=(58, 62, 70))
    draw.rectangle([0, int(h * 0.68), w, int(h * 0.72)], fill=(235, 120, 45))

    wall_count = 7
    for i in range(wall_count):
        x = 40 + i * 165 + (seed * 11) % 25
        top = 90 + (i % 4) * 25
        wall_w = 130
        base = int(h * 0.68)
        draw.polygon(
            [(x, base), (x + wall_w, base), (x + wall_w - 15, top), (x + 15, top)],
            fill=(190, 95, 55) if i % 2 else (210, 110, 60),
        )
        for row in range(9):
            for col in range(4):
                hx = x + 22 + col * 26 + (row % 2) * 8
                hy = top + 35 + row * 58
                if hy < base - 30:
                    color = (
                        (240, 240, 245),
                        (50, 130, 210),
                        (220, 60, 60),
                        (80, 180, 90),
                    )[(row + col + seed) % 4]
                    draw.ellipse([hx, hy, hx + 14, hy + 14], fill=color)

    # силуэты скалолазов
    for px, py, col in [(280, 520, (220, 220, 225)), (720, 480, (200, 205, 215)), (980, 540, (210, 215, 225))]:
        draw.ellipse([px, py - 35, px + 28, py - 5], fill=col)
        draw.line([(px + 14, py - 5), (px + 14, py + 55)], fill=col, width=6)
        draw.line([(px + 14, py + 10), (px - 20, py - 10)], fill=col, width=5)
        draw.line([(px + 14, py + 10), (px + 45, py + 5)], fill=col, width=5)

    draw.rectangle([0, 0, w, 55], fill=(20, 24, 32))
    draw.text((36, 14), title, fill=(255, 255, 255), font=_font(34, True))
    return _save(img, f"climbing_{seed}.jpg")


def make_warehouse_photo(seed: int) -> Path:
    w, h = 1200, 800
    img = Image.new("RGB", (w, h), (235, 235, 238))
    draw = ImageDraw.Draw(img)
    draw.rectangle([0, 380, w, h], fill=(210, 210, 215))
    for x in range(0, w, 180):
        draw.line([(x, 0), (x, 380)], fill=(200, 200, 205), width=3)
    draw.rectangle([120, 80, w - 120, 360], outline=(160, 160, 170), width=4)
    draw.rectangle([200 + seed * 30, 140, w - 200, 320], fill=(248, 248, 250))
    draw.text((140, 400), "Производственно-складское помещение", fill=(60, 60, 60), font=_font(28, True))
    draw.text((140, 450), f"{AREA_M2} м² · потолки {CEILING_M} м", fill=(100, 100, 100), font=_font(24))
    return _save(img, f"warehouse_{seed}.jpg")


def make_floor_plan() -> Path:
    w, h = 1100, 900
    img = Image.new("RGB", (w, h), (252, 252, 250))
    draw = ImageDraw.Draw(img)
    draw.rectangle([40, 40, w - 40, h - 40], outline=(80, 80, 90), width=3)

    zones = [
        (60, 60, 520, 420, "Боулдеринг", (200, 230, 200)),
        (540, 60, 1040, 420, "Высотные дорожки", (210, 220, 245)),
        (60, 440, 380, 820, "Детская зона", (245, 230, 210)),
        (400, 440, 700, 820, "Раздевалки + душ", (230, 230, 230)),
        (720, 440, 1040, 650, "Ресепшен + кафе", (240, 240, 220)),
        (720, 670, 1040, 820, "Прокат + склад", (225, 235, 225)),
    ]
    for x1, y1, x2, y2, label, color in zones:
        draw.rectangle([x1, y1, x2, y2], fill=color, outline=(60, 60, 70), width=2)
        draw.text((x1 + 12, y1 + 12), label, fill=(30, 30, 30), font=_font(22, True))

    draw.rectangle([50, 50, 1030, 830], outline=(46, 125, 50), width=8)
    draw.text((50, 10), f"Планировка {AREA_M2} м² (ориентировочная зонировка)", fill=(30, 30, 30), font=_font(20, True))
    draw.text((50, h - 35), "Кондратьевский пр-т, 2П — по данным объявления Avito", fill=(100, 100, 100), font=_font(16))
    return _save(img, "floor_plan.png")


def make_map() -> Path:
    w, h = 500, 350
    img = Image.new("RGB", (w, h), (230, 240, 230))
    draw = ImageDraw.Draw(img)
    for i in range(0, w, 40):
        draw.line([(i, 0), (i, h)], fill=(210, 220, 210))
    for j in range(0, h, 40):
        draw.line([(0, j), (w, j)], fill=(210, 220, 210))
    draw.line([(80, 50), (420, 280)], fill=(180, 190, 180), width=12)
    draw.line([(120, 300), (400, 80)], fill=(170, 180, 170), width=10)
    draw.ellipse([220, 130, 260, 170], fill=(220, 50, 50))
    draw.text((150, 200), "Кондратьевский пр-т, 2П", fill=(30, 30, 30), font=_font(18, True))
    draw.text((150, 225), "м. Площадь Ленина 16–20 мин", fill=(80, 80, 80), font=_font(14))
    return _save(img, "map.png")


def make_listing_card() -> Path:
    w, h = 320, 180
    img = Image.new("RGB", (w, h), (255, 255, 255))
    draw = ImageDraw.Draw(img)
    draw.rectangle([0, 0, w - 1, h - 1], outline=(200, 200, 200), width=2)
    draw.text((12, 12), "Avito · № 7985954480", fill=(0, 120, 200), font=_font(16, True))
    draw.text((12, 42), f"{AREA_M2} м² · {RENT_MONTH:,} ₽/мес".replace(",", " "), fill=(30, 30, 30), font=_font(15))
    draw.text((12, 68), ADDRESS, fill=(60, 60, 60), font=_font(13))
    draw.text((12, 95), f"{RENT_M2} ₽/м² · залог {DEPOSIT:,} ₽".replace(",", " "), fill=(60, 60, 60), font=_font(12))
    draw.text((12, 125), f"Потолки {CEILING_M} м · 1 этаж", fill=(60, 60, 60), font=_font(12))
    return _save(img, "listing_card.png")


_WAVE_PATH: Path | None = None


def _wave_image_path() -> Path | None:
    global _WAVE_PATH
    if _WAVE_PATH and _WAVE_PATH.exists():
        return _WAVE_PATH
    if not TEMPLATE_WAVE.exists():
        return None
    src = Image.open(TEMPLATE_WAVE).convert("RGB")
    sw, sh = src.size
    crop = src.crop((0, int(sh * 0.52), sw, sh))
    ASSETS.mkdir(parents=True, exist_ok=True)
    _WAVE_PATH = ASSETS / "wave_footer.png"
    crop.save(_WAVE_PATH)
    return _WAVE_PATH


def add_wave_footer(slide, prs: Presentation) -> None:
    wave_path = _wave_image_path()
    if wave_path:
        slide.shapes.add_picture(
            str(wave_path),
            0,
            prs.slide_height - Inches(1.35),
            width=prs.slide_width,
            height=Inches(1.35),
        )
        return
    top = prs.slide_height - Inches(1.05)
    shape = slide.shapes.add_shape(1, 0, top, prs.slide_width, Inches(1.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()


def set_title(slide, text: str, size: int = 36) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.font.name = "Arial"


def add_bullets(slide, items: list[str], left=0.6, top=1.35, width=7.5, size=16) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "Arial"
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)


def add_table(slide, headers: list[str], rows: list[list[str]], left, top, col_widths, font_size=11):
    tbl = slide.shapes.add_table(
        len(rows) + 1, len(headers), left, top, sum(col_widths), Inches(0.35 * (len(rows) + 2))
    ).table
    for j, w in enumerate(col_widths):
        tbl.columns[j].width = w
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xFF, 0xE0, 0xB2)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(font_size)
            p.font.name = "Arial"
            p.alignment = PP_ALIGN.CENTER
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.name = "Arial"
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            if row[0].startswith("ИТОГО") or row[0].startswith("Итого"):
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
                for p in cell.text_frame.paragraphs:
                    p.font.bold = True
    return tbl


def _try_download(url: str, dest: Path, timeout: float = 8.0) -> bool:
    import urllib.request

    try:
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            dest.write_bytes(resp.read())
        return dest.stat().st_size > 5000
    except Exception:
        return False


def build_assets() -> dict[str, Path]:
    ASSETS.mkdir(parents=True, exist_ok=True)
    remote = {
        "climb1": "https://images.unsplash.com/photo-1522163182402-834f871fd851?w=1200&q=80",
        "climb2": "https://images.unsplash.com/photo-1564769009770-7ad01a1c2d0b?w=1200&q=80",
        "climb3": "https://images.unsplash.com/photo-1517649763962-0c62306601b7?w=1200&q=80",
        "warehouse1": "https://images.unsplash.com/photo-1586528116311-ad8dd3c8310d?w=1200&q=80",
        "warehouse2": "https://images.unsplash.com/photo-1553413077-190dd305871c?w=1200&q=80",
    }
    paths: dict[str, Path] = {}
    for key, url in remote.items():
        dest = ASSETS / f"{key}.jpg"
        if not _try_download(url, dest):
            if key.startswith("climb"):
                paths[key] = make_climbing_photo(int(key[-1]) - 1, "Скалодром")
            else:
                paths[key] = make_warehouse_photo(0 if key.endswith("1") else 1)
        else:
            paths[key] = dest
    paths["floor_plan"] = make_floor_plan()
    paths["map"] = make_map()
    paths["listing"] = make_listing_card()
    return paths


def slide_title(prs, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_wave_footer(slide, prs)
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(11), Inches(2.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Кафедральный контроль. Бизнес-план скалодрома"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.name = "Arial"
    auth = slide.shapes.add_textbox(Inches(7.2), Inches(5.9), Inches(5.5), Inches(0.6))
    ap = auth.text_frame.paragraphs[0]
    ap.text = "Выполнил: _____________________, гр. _____"
    ap.font.size = Pt(18)
    ap.font.name = "Arial"
    ap.alignment = PP_ALIGN.RIGHT


def slide_concept(prs, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_wave_footer(slide, prs)
    set_title(slide, "Концепция проекта")
    add_bullets(
        slide,
        [
            "Современный крытый скалодром полного цикла: боулдеринг, верхняя прямая, детская школа.",
            f"Площадь объекта — {AREA_M2} м², потолки {CEILING_M} м (достаточно для зон до 12–14 м с навесом).",
            "Целевая аудитория: любители 18–45 лет, дети 6–16, корпоративные группы, фитнес-аудитория.",
            "Форматы: разовые визиты, абонементы, секции, аренда зала, дни рождения, корпоративы.",
            "УТП: транспортная доступность (м. Площадь Ленина), кафе, прокат, безопасные auto-belay.",
        ],
        width=6.8,
    )
    slide.shapes.add_picture(str(assets["climb1"]), Inches(7.2), Inches(1.2), width=Inches(5.5))


def slide_link(prs, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_wave_footer(slide, prs)
    set_title(slide, "Ссылка на объявление")
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.25), Inches(12), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = f"• {AVITO_URL}"
    p.font.size = Pt(11)
    p.font.name = "Arial"
    p.font.color.rgb = RGBColor(0, 102, 204)
    slide.shapes.add_picture(str(assets["warehouse1"]), Inches(0.6), Inches(1.9), width=Inches(6.1))
    slide.shapes.add_picture(str(assets["warehouse2"]), Inches(6.9), Inches(1.9), width=Inches(6.1))


def slide_characteristics(prs, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_wave_footer(slide, prs)
    set_title(slide, "Характеристика сооружения")
    add_bullets(
        slide,
        [
            f"Адрес: {ADDRESS}",
            f"Площадь помещения: {AREA_M2} м²",
            "Этаж: 1, отдельный вход с улицы",
            f"Высота потолков: {CEILING_M} м (зонирование: боулдеринг 4,5 м + верхние секции до 12 м)",
            "Состояние: чистовая отделка, бетонный пол, ворота 3×4 м",
            f"Аренда: {RENT_MONTH:,} ₽/мес ({RENT_M2} ₽/м²), залог {DEPOSIT:,} ₽".replace(",", " "),
            "Коммуникации: отопление, ХВС, канализация, 80 кВт (расширяемо)",
            "Охрана 24/7, СКУД, видеонаблюдение, пожарная сигнализация",
        ],
        width=5.8,
        size=14,
    )
    slide.shapes.add_picture(str(assets["warehouse1"]), Inches(6.5), Inches(1.1), width=Inches(3.2))
    slide.shapes.add_picture(str(assets["map"]), Inches(6.5), Inches(3.5), width=Inches(3.2))
    slide.shapes.add_picture(str(assets["listing"]), Inches(9.9), Inches(3.55), width=Inches(2.5))


def slide_layout(prs, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_wave_footer(slide, prs)
    set_title(slide, "Планировка сооружения")
    slide.shapes.add_picture(str(assets["floor_plan"]), Inches(0.35), Inches(1.15), height=Inches(5.6))
    rows = [
        ["1", "Боулдеринг (основной зал)", "280", "до 70 одновременно"],
        ["2", "Высотные дорожки + auto-belay", "220", "16 линий, 48 человек/час"],
        ["3", "Детская зона и школа", "120", "4 группы по 10 детей"],
        ["4", "Женская раздевалка + душ", "55", "20–25 человек"],
        ["5", "Мужская раздевалка + душ", "45", "15–18 человек"],
        ["6", "Ресепшен + кафе + лаунж", "70", "40 посадочных мест"],
        ["7", "Прокат, инвентарная, подсобные", "40", "хранение снаряжения"],
        ["8", "Коридоры, техзоны, СКУД", "20", "—"],
        ["Итого", "Всего", str(AREA_M2), "—"],
    ]
    add_table(
        slide,
        ["№", "Зона", "Площадь (м²)", "Вместимость / назначение"],
        rows,
        Inches(5.95),
        Inches(1.35),
        [Inches(0.35), Inches(2.0), Inches(1.1), Inches(2.5)],
        font_size=10,
    )


def slide_photos(prs, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_wave_footer(slide, prs)
    slide.shapes.add_picture(str(assets["climb1"]), Inches(0.45), Inches(0.45), width=Inches(6.2))
    slide.shapes.add_picture(str(assets["climb2"]), Inches(6.85), Inches(0.45), width=Inches(6.2))


def slide_startup(prs, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_wave_footer(slide, prs)
    set_title(slide, "Единоразовые расчёты + расчёты на инвентарь", size=28)
    left_rows = [
        [f"Аренда (залог 1-й и последний месяц) {RENT_MONTH:,} × 2".replace(",", " "), f"{RENT_MONTH * 2:,}".replace(",", " ")],
        ["Адаптация помещения (перегородки, раздевалки, душевые)", "2 200 000"],
        ["Вентиляция + отопление (усиление под зал)", "450 000"],
        ["Электрика, LED-освещение спортзон", "380 000"],
        ["Подготовка основания под маты и стены", "320 000"],
        ["Ресепшен, кафе, санузлы", "280 000"],
        ["Итого по помещению", f"{STARTUP_PREMISES:,}".replace(",", " ")],
    ]
    right_rows = [
        ["Строительство боулдеринг-стены (280 м²)", "4 500 000"],
        ["Высотные секции + металлокаркас", "3 800 000"],
        ["Наборы зацепов (стартовый комплект)", "650 000"],
        ["Маты приземления, crash pads", "920 000"],
        ["Auto-belay TRUBLUE (6 шт.)", "1 680 000"],
        ["Прокат: страховки 60, обувь 100 пар", "620 000"],
        ["Верёвки, карабины, страховочные системы", "185 000"],
        ["Детское оборудование + POS/звук", "215 000"],
        ["Итого инвентарь", f"{STARTUP_INVENTORY:,}".replace(",", " ")],
    ]
    add_table(slide, ["Статья", "Сумма (руб.)"], left_rows, Inches(0.4), Inches(1.2), [Inches(3.5), Inches(1.4)], 10)
    add_table(slide, ["Статья", "Сумма (руб.)"], right_rows, Inches(6.7), Inches(1.2), [Inches(3.6), Inches(1.4)], 10)
    summary = slide.shapes.add_textbox(Inches(0.5), Inches(6.35), Inches(12), Inches(0.6))
    p = summary.text_frame.paragraphs[0]
    p.text = f"Итого стартовых затрат (помещение + инвентарь): {STARTUP_TOTAL:,} рублей.".replace(",", " ")
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.name = "Arial"


def slide_monthly(prs, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_wave_footer(slide, prs)
    set_title(slide, "Ежемесячные расходы")
    rows = [
        ["Аренда помещения", f"{RENT_MONTH:,}".replace(",", " "), "по объявлению Avito"],
        ["Коммунальные (свет, вода, отопление, вывоз)", "120 000", f"{AREA_M2} м², душевые"],
        ["Зарплата: 4 тренера + 2 админа + менеджер", "480 000", "75+75+75+75 + 45+45 + 90, налоги"],
        ["Клининг (4 раза в неделю)", "85 000", f"~100 руб/м² × {AREA_M2} м²"],
        ["Расходники (мел, мыло, перчатки, мелки)", "15 000", ""],
        ["Интернет, телефон, CRM", "8 000", ""],
        ["Реклама (таргет, Яндекс, партнёрства)", "80 000", "ключевая статья первые 6 мес"],
        ["Страхование + обслуживание стен", "43 000", "ОСО, ТО зацепов"],
        ["ИТОГО ежемесячных расходов", f"{MONTHLY_TOTAL:,}".replace(",", " "), ""],
    ]
    add_table(
        slide,
        ["Статья", "Сумма (руб.)", "Примечание"],
        rows,
        Inches(0.5),
        Inches(1.25),
        [Inches(3.8), Inches(1.5), Inches(3.5)],
        11,
    )


def slide_revenue(prs, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_wave_footer(slide, prs)
    set_title(slide, "Доходы (прогноз 1-го месяца работы)")
    rows = [
        ["Разовые визиты (800 ₽ × ~35/день × 26)", "728 000", "консервативный трафик"],
        ["Абонементы (3 500 ₽ × 70 чел.)", "245 000", "удержание 2–3 мес"],
        ["Детские секции и ДР", "180 000", "4 группы + 6 праздников"],
        ["Корпоративы / аренда зала", "120 000", "2–3 события"],
        ["Прокат, магазин, кафе", "289 000", "маржа 35–45%"],
        ["ИТОГО доходов", f"{MONTHLY_INCOME:,}".replace(",", " "), "при активном маркетинге"],
    ]
    add_table(
        slide,
        ["Статья", "Сумма (руб.)", "Примечание"],
        rows,
        Inches(0.5),
        Inches(1.25),
        [Inches(4.2), Inches(1.5), Inches(3.2)],
        11,
    )
    note = slide.shapes.add_textbox(Inches(0.55), Inches(5.5), Inches(12), Inches(1))
    p = note.text_frame.paragraphs[0]
    p.text = (
        f"Маржинальная прибыль 1-го месяца: {MONTHLY_INCOME - MONTHLY_TOTAL:,} ₽. "
        "При росте трафика до 50–60 визитов/день выход на операционную окупаемость — 2–3-й месяц."
    ).replace(",", " ")
    p.font.size = Pt(13)
    p.font.name = "Arial"


def slide_payback(prs, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_wave_footer(slide, prs)
    set_title(slide, "Окупаемость")
    margin = MONTHLY_INCOME - MONTHLY_TOTAL
    add_bullets(
        slide,
        [
            f"Стартовые вложения: {STARTUP_TOTAL:,} руб.".replace(",", " "),
            f"Расходы в первый месяц: {MONTHLY_TOTAL:,} руб.".replace(",", " "),
            f"Доходы в первый месяц: {MONTHLY_INCOME:,} руб.".replace(",", " "),
            (
                "Точка безубыточности по операционным расходам — с 2–3-го месяца при плановом трафике. "
                "Полная окупаемость инвестиций — ориентир 20–24 месяца; консервативно 28–32 месяца."
            ),
            "Ключевой риск: длительный набор клиентской базы и сезонность спроса.",
            "Рекомендация: пробные занятия, партнёрства со школами и фитнес-клубами, корпоративные пакеты.",
        ],
        width=7.2,
        size=15,
    )
    # декоративные «иконки»
    for i, sym in enumerate(["▲", "♥"]):
        sh = slide.shapes.add_shape(1, Inches(8.5 + i * 1.2), Inches(2.5), Inches(1), Inches(1))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(0xFF, 0xC1, 0x07)
        sh.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(8.65 + i * 1.2), Inches(2.65), Inches(0.8), Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        p.text = sym
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.CENTER


def slide_thanks(prs, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_wave_footer(slide, prs)
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(10), Inches(1.5))
    p = box.text_frame.paragraphs[0]
    p.text = "Спасибо за внимание!"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.name = "Arial"


def main() -> None:
    print("Генерация изображений...")
    assets = build_assets()
    print("Сборка презентации...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs, assets)
    slide_concept(prs, assets)
    slide_link(prs, assets)
    slide_characteristics(prs, assets)
    slide_layout(prs, assets)
    slide_photos(prs, assets)
    slide_startup(prs, assets)
    slide_monthly(prs, assets)
    slide_revenue(prs, assets)
    slide_payback(prs, assets)
    slide_thanks(prs, assets)

    prs.save(OUTPUT)
    print(f"Готово: {OUTPUT}")


if __name__ == "__main__":
    main()
